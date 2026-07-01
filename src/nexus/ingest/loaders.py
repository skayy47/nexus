"""Document loaders — PDF, DOCX, XLSX, PPTX, CSV, RTF, JSON, EML, Markdown, TXT, HTML."""

from __future__ import annotations

import contextlib
import io
import re
from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)

# ── Prompt Injection Detection ───────────────────────────────────────────────
_INJECTION_PATTERNS = re.compile(
    r"(ignore\s+(all\s+)?previous\s+instructions"
    r"|you\s+are\s+now\s+a"
    r"|system\s*:\s*"
    r"|<\s*system\s*>"
    r"|IMPORTANT:\s*disregard)",
    re.IGNORECASE,
)


def sniff_injection(text: str) -> bool:
    """Check text for common prompt injection patterns."""
    return bool(_INJECTION_PATTERNS.search(text))


# ── Individual Loaders ───────────────────────────────────────────────────────


_MIN_PAGE_CHARS = 150  # Pages shorter than this are caption/figure-only — merged into neighbors


def load_pdf(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load PDF using pypdf. Falls back to Unstructured for scanned/complex PDFs.

    Pages with fewer than _MIN_PAGE_CHARS of text (e.g. figure captions) are merged
    into the nearest substantive page rather than becoming isolated low-quality chunks.
    """
    if not content:
        raise ValueError("PDF file is empty.")

    from pypdf import PdfReader

    # Attempt pypdf — catches corrupt headers, stream errors, unusual encodings.
    try:
        reader = PdfReader(io.BytesIO(content))

        if reader.is_encrypted:
            # Try blank password (some "protected" PDFs use an empty owner password)
            with contextlib.suppress(Exception):
                reader.decrypt("")
            if reader.is_encrypted:
                raise ValueError(
                    "This PDF is password-protected. Remove the password and upload again."
                )

        raw: list[tuple[int, str]] = []
        for i, page in enumerate(reader.pages, 1):
            try:
                text = (page.extract_text() or "").strip()
            except Exception:
                text = ""
            if text:
                raw.append((i, text))

    except ValueError:
        raise
    except Exception as e:
        logger.warning(
            "pypdf failed, trying Unstructured fallback",
            filename=filename,
            error=str(e),
        )
        return _load_with_unstructured(content, filename)

    if not raw:
        logger.info("pypdf extracted no text, falling back to Unstructured", filename=filename)
        return _load_with_unstructured(content, filename)

    # Merge stub pages (figure captions, blank pages) into the previous substantive page
    pages: list[dict[str, str | int]] = []
    pending_text = ""
    pending_page = raw[0][0]

    for page_num, text in raw:
        if len(text) < _MIN_PAGE_CHARS:
            # Accumulate short page text to append to the next substantive page
            pending_text += (" " + text) if pending_text else text
        else:
            combined = (pending_text + " " + text).strip() if pending_text else text
            pages.append({"text": combined, "page": page_num})
            pending_text = ""
            pending_page = page_num

    # Flush any trailing stub pages into the last substantive page
    if pending_text and pages:
        pages[-1]["text"] = (str(pages[-1]["text"]) + " " + pending_text).strip()
    elif pending_text:
        pages.append({"text": pending_text, "page": pending_page})

    logger.info(
        "PDF loaded",
        filename=filename,
        raw_pages=len(raw),
        substantive_pages=len(pages),
        merged_stubs=len(raw) - len(pages),
    )
    return pages


def _load_with_unstructured(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Fallback loader using Unstructured for scanned/complex files."""
    try:
        from unstructured.partition.auto import partition
    except ImportError as exc:
        raise ValueError(
            "Advanced PDF parsing is unavailable. Try converting to DOCX or TXT and uploading again."
        ) from exc

    try:
        elements = partition(file=io.BytesIO(content), content_type=_mime_type(filename))
    except Exception as exc:
        logger.error("Unstructured partition failed", filename=filename, error=str(exc))
        raise ValueError(
            "Could not extract text from this PDF. It may be image-only, corrupted, or use "
            "unsupported encoding. Try converting to DOCX or TXT."
        ) from exc

    pages: list[dict[str, str | int]] = []

    for element in elements:
        text = str(element).strip()
        if not text:
            continue
        page = getattr(element.metadata, "page_number", 1) or 1
        pages.append({"text": text, "page": page})

    return pages


def load_docx(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load DOCX — paragraphs + table cells for full coverage."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    entries: list[dict[str, str | int]] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            entries.append({"text": text, "page": 1})

    # Extract table content (often contains critical data)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            entries.append({"text": "\n".join(rows), "page": 1})

    return entries


def load_xlsx(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load Excel workbook — each sheet becomes a page."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    pages: list[dict[str, str | int]] = []

    for sheet_idx, sheet in enumerate(wb.worksheets, 1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            text = f"[Sheet: {sheet.title}]\n" + "\n".join(rows)
            pages.append({"text": text, "page": sheet_idx})

    return pages


def load_pptx(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load PowerPoint — each slide becomes a page, extracting text frames + notes."""
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401

    prs = Presentation(io.BytesIO(content))
    pages: list[dict[str, str | int]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []

        # Text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        if parts:
            pages.append({"text": "\n".join(parts), "page": slide_num})

    return pages


def load_csv(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load CSV — batch rows into chunks of 100 for manageable context."""
    import csv

    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = [row for row in reader if any(cell.strip() for cell in row)]

    if not rows:
        return []

    # Treat header + batches of 100 rows as separate pages
    header = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else rows
    pages: list[dict[str, str | int]] = []
    batch_size = 100

    for page_num, i in enumerate(range(0, max(len(data_rows), 1), batch_size), 1):
        batch = data_rows[i : i + batch_size]
        lines = [" | ".join(header)] + [" | ".join(r) for r in batch]
        pages.append({"text": "\n".join(lines), "page": page_num})

    return pages


def load_rtf(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load RTF using striprtf. Falls back to raw text stripping."""
    try:
        from striprtf.striprtf import rtf_to_text

        text = rtf_to_text(content.decode("latin-1", errors="replace"))
    except Exception:
        # Crude fallback: strip RTF control words
        raw = content.decode("latin-1", errors="replace")
        text = re.sub(r"\\[a-z]+\d*\s?", " ", raw)
        text = re.sub(r"[{}]", "", text)
        text = re.sub(r"\s+", " ", text).strip()

    return [{"text": text, "page": 1}] if text.strip() else []


def load_json(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load JSON — pretty-print nested structure as readable text."""
    import json

    try:
        data = json.loads(content.decode("utf-8", errors="replace"))
    except json.JSONDecodeError:
        text = content.decode("utf-8", errors="replace")
        return [{"text": text, "page": 1}] if text.strip() else []

    def flatten(obj: object, prefix: str = "") -> list[str]:
        lines: list[str] = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                key = f"{prefix}.{k}" if prefix else str(k)
                if isinstance(v, (dict, list)):
                    lines.extend(flatten(v, key))
                else:
                    lines.append(f"{key}: {v}")
        elif isinstance(obj, list):
            for i, item in enumerate(obj[:200]):  # cap at 200 items
                lines.extend(flatten(item, f"{prefix}[{i}]"))
        else:
            lines.append(f"{prefix}: {obj}")
        return lines

    text = "\n".join(flatten(data))
    return [{"text": text, "page": 1}] if text.strip() else []


def load_eml(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load EML email files — subject, sender, body (plain text preferred)."""
    import email
    from email import policy

    msg = email.message_from_bytes(content, policy=policy.default)

    parts: list[str] = []

    subject = msg.get("Subject", "")
    sender = msg.get("From", "")
    date = msg.get("Date", "")
    to = msg.get("To", "")

    if subject:
        parts.append(f"Subject: {subject}")
    if sender:
        parts.append(f"From: {sender}")
    if to:
        parts.append(f"To: {to}")
    if date:
        parts.append(f"Date: {date}")

    # Prefer plain text body; fall back to HTML stripped
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                body = part.get_content()
                break
            if ct == "text/html" and not body:
                html = part.get_content()
                body = re.sub(r"<[^>]+>", " ", html)
                body = re.sub(r"\s+", " ", body).strip()
    else:
        body = msg.get_content()

    if body:
        parts.append(f"\n{body.strip()}")

    text = "\n".join(parts)
    return [{"text": text, "page": 1}] if text.strip() else []


def load_markdown(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load Markdown file as plain text."""
    text = content.decode("utf-8", errors="replace")
    return [{"text": text, "page": 1}] if text.strip() else []


def load_txt(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load plain text file."""
    text = content.decode("utf-8", errors="replace")
    return [{"text": text, "page": 1}] if text.strip() else []


def load_html(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load HTML file, stripping tags to extract text."""
    text = content.decode("utf-8", errors="replace")
    clean = re.sub(r"<[^>]+>", " ", text)
    clean = re.sub(r"\s+", " ", clean).strip()
    return [{"text": clean, "page": 1}] if clean else []


# ── Unified Loader ──────────────────────────────────────────────────────────

_LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".doc": load_docx,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".pptx": load_pptx,
    ".ppt": load_pptx,
    ".csv": load_csv,
    ".rtf": load_rtf,
    ".json": load_json,
    ".eml": load_eml,
    ".md": load_markdown,
    ".markdown": load_markdown,
    ".txt": load_txt,
    ".html": load_html,
    ".htm": load_html,
}

SUPPORTED_EXTENSIONS = sorted(_LOADER_MAP.keys())


def load_document(content: bytes, filename: str) -> list[dict[str, str | int]]:
    """Load a document by extension. Returns [{text, page}] dicts."""
    ext = Path(filename).suffix.lower()
    loader = _LOADER_MAP.get(ext)

    if loader is None:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}"
        )

    pages = loader(content, filename)

    for page_data in pages:
        if sniff_injection(str(page_data["text"])):
            logger.warning(
                "Potential prompt injection detected",
                filename=filename,
                page=page_data.get("page"),
            )

    return pages


def _mime_type(filename: str) -> str:
    """Map filename to MIME type for Unstructured fallback."""
    ext = Path(filename).suffix.lower()
    mime_map = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".html": "text/html",
        ".htm": "text/html",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".csv": "text/csv",
        ".json": "application/json",
        ".eml": "message/rfc822",
        ".rtf": "application/rtf",
    }
    return mime_map.get(ext, "application/octet-stream")
